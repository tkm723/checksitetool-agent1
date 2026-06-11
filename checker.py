#!/usr/bin/env python3
# coding: utf-8
"""
医療・クリニック系Webサイト検品ツール（統合版）

通常実行:
  python inspect.py --existing https://www.example.com --target https://staging.example.com

LLMエージェント強化（意味的チェック追加）:
  python inspect.py --existing ... --target ... --agent
  ※ ANTHROPIC_API_KEY 環境変数が必要
"""

import argparse
import json
import re
import sys
from collections import defaultdict
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime
from urllib.parse import urljoin, urlparse

import requests
from bs4 import BeautifulSoup, Comment
import openpyxl
from openpyxl.styles import PatternFill, Font

# ── Constants ─────────────────────────────────────────────────────────────────
# agent_prompts.md の追加キーワードを反映

DUMMY_EXACT = [
    "テキストテキスト", "テキストが入ります", "テキストが入る",
    "本文が入ります", "見出しが入ります", "タイトルが入ります",
    "コンテンツが入ります", "内容が入ります", "文章が入ります",
    "説明が入ります", "キャッチが入ります", "サンプルテキスト",
    "ダミーテキスト", "テストテキスト", "テスト文章",
    "Lorem ipsum", "NOW PRINTING", "NO IMAGE", "COMING SOON",
]

DUMMY_PATTERNS = [
    re.compile(r'[〇○●□]{3,}'),               # 丸・四角3文字以上
    re.compile(r'[〇○●]+[円回ヶ月]'),           # 丸文字+単位
    re.compile(r'[〇○●]+(?:クリニック|医院|大学)'),  # 丸文字施設名
    re.compile(r'山田太郎|田中花子|田中太郎|John Doe|鈴木一郎|佐藤花子'),
]

AD_VIOLATIONS = [
    "日本一", "No.1", "NO.1", "最高", "最先端", "最新鋭", "唯一", "他院では",
    "絶対安全", "100%", "必ず", "確実に", "副作用なし", "痛みゼロ",
    "を保証", "効果を約束", "治ります", "完治します",
]

FREE_MEDICAL_SKIP = {
    '/first', '/about', '/clinic', '/access', '/info', '/recruit',
    '/career', '/privacy', '/policy', '/contact', '/faq', '/news',
    '/blog', '/staff', '/doctor', '/standard', '/facility', '/facilities',
}

CASE_URL_RE  = re.compile(r'/case|/jirei|/before|/after|/works')
CASE_H_RE    = re.compile(r'症例|ビフォーアフター')
SKIP_IMG_RE  = re.compile(r'/header|/footer|/logo|favicon|icon|blank', re.I)
MEDIA_EXT_RE = re.compile(
    r'\.(jpg|jpeg|png|gif|webp|svg|pdf|zip|mp4|mp3|doc|docx|xls|xlsx)(\?.*)?$', re.I
)
PHONE_RE = re.compile(r'0\d{1,4}[-－\-]?\d{1,4}[-－\-]?\d{4}')
PRICE_RE = re.compile(r'\d{1,3}[,，]\d{3}円|\d+円')

SKIP_URL_PATTERNS = [
    re.compile(r'/wp-admin/'),
    re.compile(r'/page/\d+/?$'),
    re.compile(r'\d{4}/\d{2}/\d{2}/'),
    MEDIA_EXT_RE,
]

DEFAULT_MODEL = "claude-sonnet-4-6"

# ── LLM prompts（--agent モード用）───────────────────────────────────────────
# チェック11〜13（自由診療・症例・未承認）の意味的判断に特化

PROMPT_MEDICAL_SEMANTIC = """\
あなたは医療・クリニック系Webサイトの自由診療・症例・未承認医薬品に関する
必要記載事項をチェックする専門エージェントです。

## 入力データ
- page: ページの抽出情報（JSON）。path・body_text・headings等が含まれる。

## チェック1: 自由診療の必要記載（料金記載があるページ）
以下のpathはスキップ：
/first /about /clinic /access /info /recruit /career /privacy /policy
/contact /faq /news /blog /staff /doctor /standard /facility /facilities

上記以外で料金（数字＋円）の記載がある場合、以下が実質的に記載されているか確認：
- 治療期間または回数（「3ヶ月後」「週1回」等の言い回しも含む）
- リスク・副作用・デメリット・注意事項（いずれか）

※ 「期間」「回数」という単語がなくても意味的に記載があればOK

## チェック2: 症例ページの必要記載
URLに /case /jirei /before /after /works を含む、
またはh1〜h3に「症例」「ビフォーアフター」を含むページで以下が揃っているか：
- 治療内容（施術内容・治療方法等も含む）
- 治療期間または回数
- 費用（料金・価格・円表記等も含む）
- リスク・副作用（デメリット・注意事項も含む）

## チェック3: 未承認医薬品の必要記載
「未承認」「個人輸入」のキーワードがあるページで以下5項目が揃っているか：
①未承認医薬品・医療機器である旨
②入手経路（個人輸入等も含む）
③国内の同一成分を含む承認医薬品等の有無
④諸外国における安全性情報
⑤医薬品副作用被害救済制度の対象外である旨

## 出力形式
JSON配列のみ返す。説明文・マークダウン不要。
[{"type": "ng", "category": "自由診療必要記載"/"症例ページ必要記載"/"未承認医薬品必要記載",
  "text": "...", "found": "...", "page_url": "...", "page_path": "..."}]
問題がなければ [] を返す。\
"""

PROMPT_AD_SEMANTIC = """\
あなたは医療広告ガイドライン違反とテキスト品質を検出する専門エージェントです。

## 入力データ
- page: ページの抽出情報（JSON）。path・body_text・headings等が含まれる。

## チェック対象外
pathが /recruit /career /privacy /policy /access で始まるページはスキップ。

## チェック1: 医療広告ガイドライン違反（グレーゾーンも検出）
明示的な違反（確実にNG）：
「日本一」「No.1」「最高」「最先端」「最新鋭」「唯一」「他院では」
「絶対安全」「100%」「必ず」「確実に」「副作用なし」「痛みゼロ」「〜を保証」「完治します」

グレーゾーン（要手動確認）：
- 「〜のような気がします」等、根拠のない効果の示唆
- 実績・件数の記載があるが対象期間が不明
- 「〜円から」「〜円〜」等の上限なし料金表示（料金記載ページのみ）
- 「税抜」「税別」「+税」等の税抜き表示

## チェック2: 患者体験談・口コミ形式（原則禁止）
「〜でよかった」「先生のおかげで」「〜でした（患者の感想）」等の
患者の声・体験談に見える表現を検出。

## チェック3: テキスト品質
- 明らかな誤字・脱字
- 表記ゆれ（「ホワイトニング」と「ホワイニング」等）
- 文中の不自然なスペースや記号

## 判定ルール
- 確実に違反 → type="ng"
- グレーゾーン・要確認 → type="manual"
- 誤検知を避けるため、文脈を読んで判断する

## 出力形式
JSON配列のみ返す。説明文・マークダウン不要。
[{"type": "ng"/"manual", "category": "医療広告GL"/"テキスト品質",
  "text": "...", "found": "...", "page_url": "...", "page_path": "..."}]
問題がなければ [] を返す。\
"""

PROMPT_BASIC_SEMANTIC = """\
あなたはWebサイトの基本情報をPython検出の補完として照合するエージェントです。
Pythonで検出が難しい院長名・診療時間・休診日・最寄り駅の比較を担当します。

## 入力データ
- page: 構築サイトのトップページ抽出情報（JSON）
- master: 既存サイトのトップページ抽出情報（JSON）

## チェック項目（Python検出済みのTEL/コピーライト/ナビ一致は除外）

### 院長名の一致
masterのbody_textから院長・ドクター名を抽出し、pageに存在するか確認。
スペース除去・部分一致・読み仮名の表記揺れを考慮。

### 診療時間・休診日の一致
masterとpageのbody_text・tableからそれぞれ診療時間と休診日を抽出して比較。
「9:30」と「9時30分」は同じと判断。曜日の略称（月/月曜/月曜日）も同じと判断。

### 最寄り駅・所要時間の一致
masterとpageで駅名・徒歩分数が一致するか確認。

### ダミー画像URL（imagesフィールド）
src に「pixta」「placeholder」「dummy」「sample」「noimage」等を含む画像を指摘。

## 判定ルール
- 誤検知を避けるため、確実に不一致と判断できる場合のみ指摘する
- 情報が見つからない場合は「要手動確認」で報告する

## 出力形式
JSON配列のみ返す。説明文・マークダウン不要。
[{"type": "ng"/"manual", "category": "基本情報照合"/"ダミー画像",
  "text": "...", "existing": "...", "found": "...", "page_url": "...", "page_path": "..."}]
問題がなければ [] を返す。\
"""

# ── URL helpers ───────────────────────────────────────────────────────────────

def get_staging_base(target_url: str) -> str:
    path = urlparse(target_url).path.strip('/')
    parts = path.split('/')
    if len(parts) >= 2 and '.' in parts[1]:
        return '/' + '/'.join(parts[:2])
    return ''


def rel_path(url: str, base_path: str) -> str:
    path = urlparse(url).path
    if base_path and path.startswith(base_path):
        path = path[len(base_path):]
    return path or '/'


def is_skip_url(url: str, base_path: str) -> bool:
    path = rel_path(url, base_path)
    for pat in SKIP_URL_PATTERNS:
        if pat.search(url):
            return True
    if re.search(r'/(blog|news|case)/?$', path):
        return True
    if re.search(r'/(blog|news)/[^/]+/?$', path):
        return True
    if 'sitemap' in url.lower():
        return True
    return False


def normalize_phone(s: str) -> str:
    return re.sub(r'[-－\-\s]', '', s)


def is_top(path: str) -> bool:
    return path in ('/', '')

# ── Crawler ───────────────────────────────────────────────────────────────────

class Crawler:
    def __init__(self, target_url: str, base_path: str):
        self.start    = target_url.rstrip('/')
        self.base_path = base_path
        self.netloc   = urlparse(target_url).netloc
        self.visited: set = set()
        self.queue: list  = [target_url]
        self.pages: dict  = {}
        self.session = requests.Session()
        self.session.headers['User-Agent'] = 'Mozilla/5.0 (compatible; WebInspector/1.0)'

    def same_domain(self, url: str) -> bool:
        return urlparse(url).netloc == self.netloc

    def crawl(self):
        while self.queue:
            url = self.queue.pop(0)
            if url in self.visited:
                continue
            if not self.same_domain(url):
                continue
            if is_skip_url(url, self.base_path):
                self.visited.add(url)
                continue
            self.visited.add(url)
            try:
                resp = self.session.get(url, timeout=15, allow_redirects=True)
                if 'html' not in resp.headers.get('content-type', ''):
                    continue
                soup = BeautifulSoup(resp.text, 'lxml')
                self.pages[url] = {'soup': soup, 'status': resp.status_code}
                for a in soup.find_all('a', href=True):
                    href = a['href'].strip()
                    if not href or href.startswith(('mailto:', 'tel:', 'javascript:')):
                        continue
                    abs_url = urljoin(url, href).split('#')[0].split('?')[0]
                    if abs_url not in self.visited and self.same_domain(abs_url):
                        self.queue.append(abs_url)
            except Exception as e:
                print(f'  エラー: {url} — {e}', file=sys.stderr)
                self.pages[url] = {'soup': None, 'status': 0}

# ── Result builders ───────────────────────────────────────────────────────────

def ng(category, content, location='', existing='', detected=''):
    return dict(kind='NG', category=category, content=content,
                location=location, existing=existing, detected=detected)

def warn(category, content, location=''):
    return dict(kind='要手動確認', category=category, content=content,
                location=location, existing='', detected='')

# ── Check 1: Dummy text（DOMノードベース・位置特定・×N件数）────────────────────

def _heading_context(node) -> str:
    ancestors = []
    p = node.parent
    while p and p.name not in ('[document]', 'html', 'body'):
        ancestors.append(p)
        p = p.parent

    for anc in ancestors:
        if anc.name == 'table':
            row_label = ''
            for a in ancestors:
                if a.name in ('th', 'td'):
                    row = a.find_parent('tr')
                    if row:
                        for th in row.find_all('th'):
                            t = th.get_text(strip=True)
                            if t:
                                row_label = t
                                break
                    break
            table_h = ''
            prev = anc
            while prev:
                prev = prev.find_previous_sibling()
                if prev and hasattr(prev, 'name') and prev.name in ('h1','h2','h3','h4'):
                    table_h = prev.get_text(strip=True)
                    break
            if table_h and row_label:
                return f'h4:「{table_h}」> 表の行「{row_label}」'
            if table_h:
                return f'h4:「{table_h}」> 表内'
            break

    for anc in ancestors:
        if anc.name in ('dd', 'dt'):
            dl = anc.find_parent('dl')
            if dl:
                for dt in dl.find_all('dt'):
                    t = dt.get_text(strip=True)
                    if t:
                        return f'項目「{t}」'
            break

    for anc in ancestors:
        if anc.name in ('h1','h2','h3','h4'):
            return f'{anc.name}:「{anc.get_text(strip=True)}」'

    parent = node.parent
    if parent:
        sib = parent
        while sib:
            sib = sib.find_previous_sibling()
            if sib and hasattr(sib, 'name') and sib.name in ('h1','h2','h3','h4'):
                return f'{sib.name}:「{sib.get_text(strip=True)}」'
    return ''


def check_dummy_text(soup: BeautifulSoup) -> list:
    found: dict = defaultdict(int)
    for node in soup.find_all(string=True):
        if isinstance(node, Comment):
            continue
        text = node.strip()
        if not text:
            continue
        matched = None
        for kw in DUMMY_EXACT:
            if kw.lower() in text.lower():
                matched = kw
                break
        if not matched:
            for pat in DUMMY_PATTERNS:
                m = pat.search(text)
                if m:
                    matched = m.group()
                    break
        if matched:
            found[(matched, _heading_context(node))] += 1

    results = []
    for (kw, heading), count in found.items():
        suffix = f' ×{count}箇所' if count > 1 else ''
        results.append(ng('ダミーテキスト', f'ダミーテキスト検出: 「{kw}」{suffix}', heading))
    return results

# ── Check 2: External links without target="_blank" ───────────────────────────

def check_external_links(soup: BeautifulSoup, page_url: str) -> list:
    own  = urlparse(page_url).netloc
    seen = set()
    results = []
    for a in soup.find_all('a', href=True):
        href = a['href']
        if not href or href.startswith(('#', 'mailto:', 'tel:', 'javascript:')):
            continue
        try:
            domain = urlparse(urljoin(page_url, href)).netloc
        except Exception:
            continue
        if domain and domain != own and a.get('target') != '_blank':
            key = href[:80]
            if key not in seen:
                seen.add(key)
                results.append(ng('外部リンク', f'target="_blank" なし: 「{a.get_text(strip=True)[:40]}」',
                                  f'href="{href}"'))
    return results

# ── Check 3: Empty / invalid links ───────────────────────────────────────────

def check_empty_links(soup: BeautifulSoup) -> list:
    seen = set()
    results = []
    for a in soup.find_all('a', href=True):
        if a.find_parent('nav'):
            continue
        href = a['href'].strip()
        if href.startswith('#') and len(href) > 1:
            continue
        if href in ('', '#', 'javascript:void(0)', 'javascript:void(0);', 'javascript:;'):
            text = a.get_text(strip=True)[:40] or '(テキストなし)'
            key  = (text, href)
            if key not in seen:
                seen.add(key)
                results.append(ng('空リンク', f'無効リンク: 「{text}」 href="{href}"'))
    return results

# ── Check 4: Missing alt ──────────────────────────────────────────────────────

def check_missing_alt(soup: BeautifulSoup) -> list:
    results = []
    MEANINGLESS_ALT = re.compile(r'^(image|img|photo|写真|画像|pic)$', re.I)
    for img in soup.find_all('img'):
        src = img.get('src', '')
        if SKIP_IMG_RE.search(src):
            continue
        if not img.has_attr('alt'):
            results.append(ng('alt未設定', f'alt属性未設定: src="{src[:60]}"'))
        elif MEANINGLESS_ALT.match(img['alt'].strip()):
            results.append(ng('alt未設定', f'altが無意味な値: alt="{img["alt"]}" src="{src[:50]}"'))
    return results

# ── Check 5: h1 tag ───────────────────────────────────────────────────────────

def check_h1(soup: BeautifulSoup) -> list:
    h1s = soup.find_all('h1')
    if not h1s:
        return [ng('h1タグ', 'h1タグが存在しない')]
    if len(h1s) >= 2:
        texts = [h.get_text(strip=True)[:30] for h in h1s]
        return [ng('h1タグ', f'h1タグが{len(h1s)}個ある: {texts}')]
    # h1にダミーワードが含まれていないか
    h1_text = h1s[0].get_text(strip=True)
    if re.search(r'テスト|テキスト|ダミー|サンプル', h1_text):
        return [ng('h1タグ', f'h1にダミー語が含まれている: 「{h1_text[:40]}」')]
    return []

# ── Check 6: Heading hierarchy ────────────────────────────────────────────────

def check_heading_hierarchy(soup: BeautifulSoup) -> list:
    headings = soup.find_all(['h1','h2','h3','h4','h5','h6'])
    levels   = [int(h.name[1]) for h in headings]
    results  = []
    for i in range(1, len(levels)):
        if levels[i] > levels[i-1] + 1:
            results.append(ng('見出し階層',
                f'見出し階層スキップ: h{levels[i-1]}→h{levels[i]} 「{headings[i].get_text(strip=True)[:30]}」'))
    return results

# ── Check 7: Phone number hyphen consistency ──────────────────────────────────

def check_phone_hyphen(soup: BeautifulSoup) -> list:
    # tel: href はディスプレイ値でないため先に除去
    tel_tags = list(soup.find_all('a', href=re.compile(r'^tel:')))
    for a in tel_tags:
        a.decompose()

    phones = PHONE_RE.findall(soup.get_text())
    if len(phones) < 2:
        return []
    has_hyp = any(re.search(r'[-－\-]', p) for p in phones)
    no_hyp  = any(not re.search(r'[-－\-]', p) for p in phones)
    if has_hyp and no_hyp:
        return [ng('電話番号', f'ハイフンあり・なし混在: {list(dict.fromkeys(phones))[:6]}')]
    return []

# ── Check 8: Breadcrumb ───────────────────────────────────────────────────────

def check_breadcrumb(soup: BeautifulSoup, path: str) -> list:
    if is_top(path):
        return []
    bc = (
        soup.find(class_=re.compile(r'breadcrumb', re.I)) or
        soup.find(id=re.compile(r'breadcrumb', re.I)) or
        soup.find('nav', {'aria-label': re.compile(r'breadcrumb|パンくず', re.I)}) or
        soup.find(class_=re.compile(r'\bpan\b|pankuzu', re.I))
    )
    if not bc:
        return [warn('パンくず', 'パンくずが見つからない（要手動確認）')]
    items = bc.find_all(['li', 'span', 'a'])
    if items:
        last = items[-1].get_text(strip=True)
        h1   = soup.find('h1')
        if h1 and last and last != h1.get_text(strip=True):
            return [ng('パンくず',
                f'パンくず末尾「{last}」とh1「{h1.get_text(strip=True)}」が不一致')]
    return []

# ── Check 9: reCAPTCHA ────────────────────────────────────────────────────────

def check_recaptcha(soup: BeautifulSoup) -> list:
    if not soup.find('form'):
        return []
    html = str(soup)
    if 'recaptcha' not in html.lower() and 'grecaptcha' not in html:
        return [ng('reCAPTCHA', 'フォームあり・reCAPTCHAなし')]
    return []

# ── Check 10: Price dummy ─────────────────────────────────────────────────────

def check_price_dummy(soup: BeautifulSoup, path: str) -> list:
    if re.search(r'/recruit|/career', path):
        return []
    results = []
    text = soup.get_text()

    for m in re.finditer(r'[0０]{2,}[,，][0０]{3}円', text):
        results.append(ng('料金ダミー', f'全桁ゼロの金額: 「{m.group()}」'))
    for m in re.finditer(r'[〇○●]{2,}円', text):
        results.append(ng('料金ダミー', f'丸文字の金額: 「{m.group()}」'))
    for m in re.finditer(
        r'[\d,，]+円[〜～][0０]{2,}[,，][0０]{3}円|[0０]{2,}[,，][0０]{3}円[〜～][\d,，]+円', text
    ):
        results.append(ng('料金ダミー', f'料金範囲の片方がゼロ: 「{m.group()}」'))
    if PRICE_RE.search(text) and not re.search(r'税込|税抜|消費税', text):
        results.append(ng('料金ダミー', '料金記載あり・「税込」表記なし'))
    for m in re.finditer(r'[〇○●]+円[〜～]\s*$', text, re.M):
        results.append(ng('料金ダミー', f'下限のみの料金表示: 「{m.group().strip()}」'))
    return results

# ── Checks 11〜14: Python fallback（--agent なし時に使用）─────────────────────

def check_free_medical_py(soup: BeautifulSoup, path: str) -> list:
    p = path.rstrip('/')
    for skip in FREE_MEDICAL_SKIP:
        s = skip.rstrip('/')
        if p == s or p.startswith(s + '/'):
            return []
    text = soup.get_text()
    if not PRICE_RE.search(text):
        return []
    results = []
    if not re.search(r'治療期間|治療回数|施術回数|期間|回数', text):
        results.append(ng('自由診療必要記載', '料金記載あり・治療期間または回数の記載なし'))
    if not re.search(r'注意事項|デメリット|副作用|リスク', text):
        results.append(ng('自由診療必要記載', '料金記載あり・注意事項/デメリット/副作用/リスクの記載なし'))
    return results


def check_case_page_py(soup: BeautifulSoup, path: str) -> list:
    if is_top(path):
        return []
    is_case = bool(CASE_URL_RE.search(path))
    if not is_case:
        for tag in soup.find_all(['h1','h2','h3']):
            if CASE_H_RE.search(tag.get_text()):
                is_case = True
                break
    if not is_case:
        return []
    text = soup.get_text()
    required = [
        ('治療内容',         r'治療内容|施術内容|治療方法'),
        ('治療期間または回数', r'治療期間|治療回数|施術回数|期間|回数'),
        ('費用',             r'費用|料金|価格|\d+円'),
        ('リスク・副作用',   r'リスク|副作用|デメリット|注意事項'),
    ]
    return [ng('症例ページ必要記載', f'症例ページ: 「{label}」の記載なし')
            for label, pat in required if not re.search(pat, text)]


def check_unapproved_drug_py(soup: BeautifulSoup) -> list:
    text = soup.get_text()
    if not any(kw in text for kw in ('未承認', '個人輸入', '自由診療')):
        return []
    required = [
        ('未承認である旨',           r'未承認'),
        ('入手経路',                 r'入手経路|個人輸入'),
        ('国内の承認医薬品等の有無', r'承認医薬品|国内承認'),
        ('諸外国の安全性情報',       r'諸外国|海外.*安全'),
        ('副作用被害救済制度の対象外', r'副作用被害救済'),
    ]
    return [ng('未承認医薬品必要記載', f'「{label}」の記載なし')
            for label, pat in required if not re.search(pat, text)]


def check_ad_guideline_py(soup: BeautifulSoup, path: str) -> list:
    for skip in ('/recruit', '/access', '/privacy', '/policy'):
        if path.startswith(skip):
            return []
    text    = soup.get_text()
    results = [ng('医療広告GL', f'違反キーワード: 「{kw}」') for kw in AD_VIOLATIONS if kw in text]
    if re.search(r'[（(][3-9]\d[代女男性さん][）)]|患者.{0,10}コメント|体験談', text):
        results.append(ng('医療広告GL', '患者体験談・口コミに見える表現を検出'))
    return results

# ── Check 15: Basic info comparison（Pythonパート）──────────────────────────

def check_basic_info_py(t_soup: BeautifulSoup, e_soup: BeautifulSoup) -> list:
    results = []
    CLINIC_KW = ['クリニック', '医院', 'デンタル', '歯科', '整形', '皮膚科', '内科']

    def phones_from(soup):
        nums = set()
        for a in soup.find_all('a', href=re.compile(r'^tel:')):
            nums.add(normalize_phone(a['href'].replace('tel:', '')))
        for m in PHONE_RE.finditer(soup.get_text()):
            nums.add(normalize_phone(m.group()))
        return nums - {''}

    # 電話番号照合
    missing = phones_from(e_soup) - phones_from(t_soup)
    if missing:
        results.append(ng('基本情報照合', f'既存サイトの電話番号が見つからない: {sorted(missing)}',
                          existing=str(sorted(phones_from(e_soup))),
                          detected=str(sorted(phones_from(t_soup)))))

    # TELリンク表示番号一致
    for a in t_soup.find_all('a', href=re.compile(r'^tel:')):
        href_n = normalize_phone(a['href'].replace('tel:', ''))
        disp_n = normalize_phone(a.get_text(strip=True))
        if href_n and disp_n and href_n != disp_n:
            results.append(ng('基本情報照合',
                f'TELリンクと表示番号が不一致: href={href_n} 表示={disp_n}'))

    # 院名照合
    def clinic_name(soup):
        footer = soup.find('footer')
        if footer:
            for s in footer.find_all(string=True):
                s = s.strip()
                if any(kw in s for kw in CLINIC_KW) and len(s) > 3:
                    return s[:40]
        title = soup.find('title')
        if title:
            for part in re.split(r'[|｜\-－]', title.get_text()):
                part = part.strip()
                if any(kw in part for kw in CLINIC_KW):
                    return part
        return ''

    t_name, e_name = clinic_name(t_soup), clinic_name(e_soup)
    if t_name and e_name and t_name not in e_name and e_name not in t_name:
        results.append(ng('基本情報照合', f'院名が不一致: 既存「{e_name}」 構築「{t_name}」',
                          existing=e_name, detected=t_name))

    # コピーライト年号
    footer = t_soup.find('footer')
    if footer:
        nav = footer.find('nav')
        ft = footer.get_text()
        if nav:
            ft = ft.replace(nav.get_text(), '')
        m = re.search(r'(?:©|Copyright)[^\d]*(\d{4})', ft)
        if m:
            year = int(m.group(1))
            if year != datetime.now().year:
                results.append(ng('基本情報照合',
                    f'コピーライト年号が現在年と不一致: {year} (現在: {datetime.now().year})'))

    # ナビテキスト一致
    nav_map: dict = {}
    for nav in t_soup.find_all('nav'):
        for a in nav.find_all('a', href=True):
            href = a['href']
            text = a.get_text(strip=True)
            if not text:
                continue
            if href in nav_map and nav_map[href] != text:
                results.append(ng('基本情報照合',
                    f'同じhrefでナビテキストが異なる: href="{href}" '
                    f'「{nav_map[href]}」と「{text}」'))
            else:
                nav_map[href] = text

    return results

# ── Check 15b: 基本情報照合（手動入力版・新規案件用）────────────────────────────

def check_basic_info_manual(t_soup: BeautifulSoup, manual_info: dict) -> list:
    """手動入力された基本情報がターゲットサイトに存在するか照合する（既存サイトなし案件用）"""
    results = []
    text = t_soup.get_text(' ', strip=True)

    # 電話番号
    if manual_info.get('tel'):
        tel_norm = normalize_phone(manual_info['tel'])
        page_phones = set()
        for a in t_soup.find_all('a', href=re.compile(r'^tel:')):
            page_phones.add(normalize_phone(a['href'].replace('tel:', '')))
        for m in PHONE_RE.finditer(text):
            page_phones.add(normalize_phone(m.group()))
        page_phones.discard('')
        if tel_norm and tel_norm not in page_phones:
            results.append(ng('基本情報照合', f'指定の電話番号が見つからない: {manual_info["tel"]}',
                              existing=manual_info['tel'],
                              detected=str(sorted(page_phones)) if page_phones else '（なし）'))

    # FAX番号
    if manual_info.get('fax'):
        fax_norm = re.sub(r'[^\d]', '', manual_info['fax'])
        if fax_norm and fax_norm not in re.sub(r'[^\d]', '', text):
            results.append(warn('基本情報照合', f'FAX番号が見つからない可能性: {manual_info["fax"]}',
                                existing=manual_info['fax']))

    # 院名
    if manual_info.get('clinic_name'):
        name = manual_info['clinic_name'].strip()
        if name and name not in text:
            results.append(ng('基本情報照合', f'院名が見つからない: 「{name}」',
                              existing=name))

    # 院長名
    if manual_info.get('director'):
        director = manual_info['director'].strip()
        if director and director not in text:
            results.append(ng('基本情報照合', f'院長名が見つからない: 「{director}」',
                              existing=director))

    # 診療時間（時刻形式の有無で確認）
    if manual_info.get('hours'):
        if not re.search(r'\d{1,2}[:：]\d{2}', text):
            results.append(warn('基本情報照合', '診療時間の記載が見つからない可能性があります',
                                existing=manual_info['hours']))

    # 休診日
    if manual_info.get('closed'):
        closed = manual_info['closed'].strip()
        if closed and closed not in text:
            results.append(warn('基本情報照合', f'休診日の記載が確認できません: 「{closed}」',
                                existing=closed))

    # 郵便番号・住所
    if manual_info.get('address'):
        address = manual_info['address'].strip()
        postal_m = re.search(r'[〒]?\d{3}[-－]?\d{4}', address)
        if postal_m:
            postal = re.sub(r'[^0-9]', '', postal_m.group())
            if postal and postal not in re.sub(r'[^0-9]', '', text):
                results.append(ng('基本情報照合', f'郵便番号が見つからない: {postal_m.group()}',
                                  existing=address))
        else:
            addr_part = re.sub(r'^〒?\d{3}[-－]?\d{4}\s*', '', address).strip()[:10]
            if addr_part and addr_part not in text:
                results.append(warn('基本情報照合', '住所の記載が確認できません',
                                    existing=address))

    # 最寄り駅
    if manual_info.get('station'):
        station = manual_info['station'].strip()
        station_m = re.search(r'[\w぀-鿿]+駅', station)
        station_name = station_m.group() if station_m else station
        if station_name and station_name not in text:
            results.append(warn('基本情報照合', f'最寄り駅の記載が確認できません: 「{station_name}」',
                                existing=station))

    return results


# ── Check 16: Meta tags ───────────────────────────────────────────────────────

def check_meta_tags(soup: BeautifulSoup) -> list:
    results = []
    title = soup.find('title')
    if not title or not title.get_text(strip=True):
        results.append(ng('metaタグ', 'titleタグが空または未設定'))
    desc = soup.find('meta', attrs={'name': re.compile(r'^description$', re.I)})
    if not desc or not (desc.get('content') or '').strip():
        results.append(ng('metaタグ', 'meta descriptionが未設定'))
    kw = soup.find('meta', attrs={'name': re.compile(r'^keywords$', re.I)})
    if not kw or not (kw.get('content') or '').strip():
        results.append(ng('metaタグ', 'meta keywordsが未設定'))
    # OGP
    for prop in ('og:title', 'og:description', 'og:image'):
        tag = soup.find('meta', property=re.compile(f'^{prop}$', re.I))
        if not tag or not (tag.get('content') or '').strip():
            results.append(ng('metaタグ', f'{prop} が未設定'))
    # noindex / canonical
    robots = soup.find('meta', attrs={'name': re.compile(r'^robots$', re.I)})
    if robots and 'noindex' in (robots.get('content') or '').lower():
        results.append(ng('metaタグ', 'noindexが設定されている（staging設定の残り？）'))
    if not soup.find('link', rel='canonical'):
        results.append(warn('metaタグ', 'canonicalタグが未設定（要確認）'))
    return results

# ── Check 17: Post migration ──────────────────────────────────────────────────

def check_post_migration(e_soup: BeautifulSoup, t_soup: BeautifulSoup) -> list:
    titles = set()
    for a in e_soup.find_all('a', href=True):
        if re.search(r'/(blog|news)/[^/]+/?$', a['href']):
            t = a.get_text(strip=True)
            if t and len(t) > 5:
                titles.add(t[:50])
    if not titles:
        return []
    t_text  = t_soup.get_text()
    missing = [t for t in list(titles)[:15] if t not in t_text]
    if missing:
        return [warn('投稿記事移行',
            f'既存サイトの記事タイトルが見つからない（要確認）: {missing[:3]}')]
    return []

# ── Manual checks ─────────────────────────────────────────────────────────────

MANUAL_ITEMS = [
    '【レイアウト】デザインカンプとの照合・レイアウト崩れ（PC/SP/iPad）・WEBフォント確認',
    '【ダミー画像】グレーの四角のみのダミー画像（テキストなし）は自動検知不可のため全ページ目視確認',
    '【投稿】テスト投稿が削除されているか・投稿一覧・アイキャッチの表示崩れ',
    '【フォーム】テスト送信・必須項目・自動返信メール・reCAPTCHA・完了ページへの遷移',
    '【その他】Googleマップのピン位置・404ページ・ページトップボタン・JSの動作・画像WebP変換',
    '【基本情報】設備・器具等の名称・最寄り駅からの所要時間がサイト内で統一されているか',
]

# ── Wireframe check ───────────────────────────────────────────────────────────

def check_wireframe(wire_path: str, pages: dict) -> list:
    try:
        from pptx import Presentation
    except ImportError:
        print('python-pptx が見つかりません。pip install python-pptx', file=sys.stderr)
        return []
    try:
        prs = Presentation(wire_path)
    except Exception as e:
        print(f'PPTXを開けません: {e}', file=sys.stderr)
        return []

    all_text = '\n'.join(d['soup'].get_text() for d in pages.values() if d.get('soup'))
    all_norm = re.sub(r'[-－\-〒\s]', '', all_text)

    DUMMY_PHONE_RE  = re.compile(r'0{3,}[-\-]?\d+[-\-]?\d+|[xX*]{3}')
    DUMMY_POSTAL_RE = re.compile(r'000-0000|[xX*]{3}-[xX*]{4}', re.I)
    PHONE_WIRE  = re.compile(r'0\d{1,4}[-－\-]?\d{1,4}[-－\-]?\d{4}')
    POSTAL_WIRE = re.compile(r'〒?\d{3}[-\-]\d{4}')
    CP_WIRE     = re.compile(r'(?:©|Copyright)[^\n]*\d{4}[^\n]*')

    results = []
    for idx, slide in enumerate(prs.slides, 1):
        slide_text = '\n'.join(s.text for s in slide.shapes if hasattr(s, 'text'))

        for m in PHONE_WIRE.finditer(slide_text):
            p = m.group()
            if not DUMMY_PHONE_RE.search(p) and normalize_phone(p) not in all_norm:
                results.append(ng('ワイヤー照合',
                    f'スライド{idx}: 電話番号がサイトに見つからない: 「{p}」'))

        for m in POSTAL_WIRE.finditer(slide_text):
            postal = m.group()
            if not DUMMY_POSTAL_RE.search(postal) and re.sub(r'[-\-〒]', '', postal) not in all_norm:
                results.append(ng('ワイヤー照合',
                    f'スライド{idx}: 郵便番号がサイトに見つからない: 「{postal}」'))

        for m in CP_WIRE.finditer(slide_text):
            cp = m.group().strip()
            if cp not in all_text:
                results.append(warn('ワイヤー照合',
                    f'スライド{idx}: コピーライトの照合が必要: 「{cp}」'))

        for line in slide_text.split('\n'):
            line = line.strip()
            if len(line) >= 100 and not any(d in line for d in DUMMY_EXACT):
                if line[:30] not in all_text:
                    results.append(warn('ワイヤー照合',
                        f'スライド{idx}: 本文テキストが未反映の可能性: 「{line[:40]}…」'))
    return results

# ── LLM helpers（--agent モード）──────────────────────────────────────────────

def _extract_page_data_for_llm(url: str, soup: BeautifulSoup, base_path: str) -> dict:
    """LLMエージェントへ渡す軽量な抽出データを生成する。"""
    path = rel_path(url, base_path)
    for tag in soup(['script', 'style']):
        tag.decompose()
    body_text = soup.get_text(separator=' ', strip=True)[:5000]
    headings  = [{'level': int(h.name[1]), 'text': h.get_text(strip=True)[:80]}
                 for h in soup.find_all(['h1','h2','h3','h4'])][:30]
    images    = [{'src': img.get('src', '')[:100], 'alt': img.get('alt')}
                 for img in soup.find_all('img')][:30]
    return {
        'url': url,
        'path': path,
        'is_top': is_top(path),
        'headings': headings,
        'images': images,
        'body_text': body_text,
    }


def _parse_llm_json(text: str) -> list:
    text = re.sub(r'^```(?:json)?\s*', '', text.strip(), flags=re.M)
    text = re.sub(r'```\s*$', '', text.strip(), flags=re.M)
    m = re.search(r'\[[\s\S]*\]', text.strip())
    try:
        return json.loads(m.group() if m else text.strip())
    except Exception:
        return []


def _run_llm_agent(client, name: str, system: str, payload: dict, model: str) -> list:
    try:
        resp = client.messages.create(
            model=model,
            max_tokens=4096,
            system=system,
            messages=[{'role': 'user', 'content': json.dumps(payload, ensure_ascii=False)}],
        )
        return _parse_llm_json(resp.content[0].text)
    except Exception as e:
        print(f'  LLMエージェント({name})エラー: {e}', file=sys.stderr)
        return []


def _llm_to_ng_warn(issue: dict) -> dict:
    """LLMの {type, category, text, found, ...} を統一形式に変換する。"""
    t = issue.get('type', 'ng')
    if t == 'ng':
        return ng(
            issue.get('category', 'LLMチェック'),
            issue.get('text', ''),
            issue.get('heading', issue.get('location', '')),
            issue.get('existing', ''),
            issue.get('found', ''),
        )
    return warn(
        issue.get('category', 'LLMチェック'),
        issue.get('text', ''),
        issue.get('heading', issue.get('location', '')),
    )


def run_llm_checks(client, page_data_llm: dict, master_data_llm: dict, model: str) -> list:
    """チェック11〜14をLLMで並列実行し、統一形式に変換して返す。"""
    tasks = [
        ('医療必要記載',   PROMPT_MEDICAL_SEMANTIC, {'page': page_data_llm, 'master': master_data_llm}),
        ('医療広告GL',     PROMPT_AD_SEMANTIC,      {'page': page_data_llm}),
    ]
    if page_data_llm.get('is_top') and master_data_llm:
        tasks.append(('基本情報照合LLM', PROMPT_BASIC_SEMANTIC,
                       {'page': page_data_llm, 'master': master_data_llm}))

    all_issues: list = []
    with ThreadPoolExecutor(max_workers=len(tasks)) as ex:
        futures = {ex.submit(_run_llm_agent, client, name, prompt, payload, model): name
                   for name, prompt, payload in tasks}
        for future in as_completed(futures):
            for issue in future.result():
                all_issues.append(_llm_to_ng_warn(issue))

    return all_issues


def _dedup_merge(py_results: list, llm_results: list) -> list:
    """Python結果を優先し、LLMが新規に見つけた指摘のみ追加する。"""
    def norm_key(r):
        # kind + category + content の先頭50文字で同一判定
        return (r.get('kind',''), r.get('category',''), r.get('content','')[:50])

    py_keys = {norm_key(r) for r in py_results}
    merged  = list(py_results)
    for r in llm_results:
        if norm_key(r) not in py_keys:
            merged.append(r)
    return merged

# ── Per-page check runner ─────────────────────────────────────────────────────

def run_checks(url: str, page_data: dict, base_path: str,
               existing_soup=None, llm_client=None, model=DEFAULT_MODEL,
               manual_info=None) -> list:
    soup = page_data.get('soup')
    if not soup:
        return [ng('クロール', f'ページ取得失敗 (status={page_data.get("status","?")})', '')]

    path    = rel_path(url, base_path)
    results = []

    # ─ Python checks 1〜10（常時実行）─────────────────────────────────────────
    results += check_dummy_text(soup)
    results += check_external_links(soup, url)
    results += check_empty_links(soup)
    results += check_missing_alt(soup)
    results += check_h1(soup)
    results += check_heading_hierarchy(soup)
    results += check_phone_hyphen(soup)
    results += check_breadcrumb(soup, path)
    results += check_recaptcha(soup)
    results += check_price_dummy(soup, path)

    # ─ チェック11〜14: Python fallback or LLM ────────────────────────────────
    if llm_client:
        page_llm   = _extract_page_data_for_llm(url, soup, base_path)
        master_llm = (_extract_page_data_for_llm(existing_soup.find('html').get('data-url', url),
                                                   existing_soup, '')
                      if existing_soup else {})
        llm_results = run_llm_checks(llm_client, page_llm, master_llm, model)

        # Python も走らせて新規検出分をLLMで補完（精度向上のため両方採用）
        py_semantic = (check_free_medical_py(soup, path)
                       + check_case_page_py(soup, path)
                       + check_unapproved_drug_py(soup)
                       + check_ad_guideline_py(soup, path))
        results += _dedup_merge(py_semantic, llm_results)
    else:
        results += check_free_medical_py(soup, path)
        results += check_case_page_py(soup, path)
        results += check_unapproved_drug_py(soup)
        results += check_ad_guideline_py(soup, path)

    # ─ チェック15: Python（TEL/コピーライト/ナビ）────────────────────────────
    if is_top(path):
        if existing_soup:
            results += check_basic_info_py(soup, existing_soup)
        elif manual_info:
            results += check_basic_info_manual(soup, manual_info)

    # ─ チェック16〜17 + 手動（トップページのみ）─────────────────────────────
    if is_top(path):
        results += check_meta_tags(soup)
        if existing_soup:
            results += check_post_migration(existing_soup, soup)
        results += [warn('手動確認', item) for item in MANUAL_ITEMS]

    return results

# ── Excel output ──────────────────────────────────────────────────────────────

def write_excel(all_results: dict, output):
    wb    = openpyxl.Workbook()
    F_NG  = PatternFill(fill_type='solid', fgColor='FFCCCC')
    F_WRN = PatternFill(fill_type='solid', fgColor='FFFACC')
    F_OK  = PatternFill(fill_type='solid', fgColor='CCFFCC')
    F_HDR = PatternFill(fill_type='solid', fgColor='4472C4')
    FHDR  = Font(color='FFFFFF', bold=True)

    def fill(kind):
        return F_NG if kind == 'NG' else F_WRN if kind == '要手動確認' else F_OK

    # Sheet 1: Summary
    ws1 = wb.active
    ws1.title = 'サマリー'
    for c, h in enumerate(['ページパス', 'NG件数', '要手動確認件数', '主なNGカテゴリ'], 1):
        ws1.cell(row=1, column=c, value=h).fill = F_HDR
        ws1.cell(row=1, column=c).font = FHDR
    for r, (url, res) in enumerate(all_results.items(), 2):
        path  = rel_path(url, '')
        ng_n  = sum(1 for x in res if x['kind'] == 'NG')
        w_n   = sum(1 for x in res if x['kind'] == '要手動確認')
        cats  = list(dict.fromkeys(x['category'] for x in res if x['kind'] == 'NG'))[:3]
        rf    = F_NG if ng_n else F_WRN if w_n else F_OK
        for c, v in enumerate([path, ng_n, w_n, ', '.join(cats)], 1):
            ws1.cell(row=r, column=c, value=v).fill = rf
    ws1.column_dimensions['A'].width = 40
    ws1.column_dimensions['B'].width = 10
    ws1.column_dimensions['C'].width = 15
    ws1.column_dimensions['D'].width = 45

    # Sheet 2: Detail
    ws2 = wb.create_sheet('ページ別詳細')
    for c, h in enumerate(
        ['ページパス','構築URL','既存URL','種別','カテゴリ','内容','既存値','検出値','箇所'], 1
    ):
        ws2.cell(row=1, column=c, value=h).fill = F_HDR
        ws2.cell(row=1, column=c).font = FHDR

    row = 2
    for url, res in all_results.items():
        path = rel_path(url, '')
        for x in res:
            f = fill(x['kind'])
            for c, v in enumerate(
                [path, url, '', x['kind'], x['category'],
                 x['content'], x.get('existing',''), x.get('detected',''), x.get('location','')], 1
            ):
                ws2.cell(row=row, column=c, value=v).fill = f
            row += 1

    import openpyxl.utils as xu
    for c, w in enumerate([35, 50, 50, 12, 20, 70, 30, 30, 35], 1):
        ws2.column_dimensions[xu.get_column_letter(c)].width = w

    wb.save(output)  # accepts file path (str) or BytesIO

# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description='医療・クリニック系Webサイト検品ツール',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
例:
  # 通常実行（Python検品のみ）
  python inspect.py --existing https://www.example.com --target https://staging.example.com

  # LLM強化（意味的チェック追加・ANTHROPIC_API_KEY 必須）
  python inspect.py --existing ... --target ... --agent

  # ワイヤーフレーム照合も含める
  python inspect.py --existing ... --target ... --wire wire.pptx --agent
        """
    )
    parser.add_argument('--existing', required=True, help='既存サイトURL（正解データ）')
    parser.add_argument('--target',   required=True, help='構築サイトURL（チェック対象）')
    parser.add_argument('--wire',     default=None,  help='ワイヤーフレームPPTXファイルパス')
    parser.add_argument('--agent',    action='store_true',
                        help='LLMエージェントによる意味的チェックを追加（ANTHROPIC_API_KEY 必須）')
    parser.add_argument('--model',    default=DEFAULT_MODEL,
                        help=f'LLMモデルID (default: {DEFAULT_MODEL})')
    args = parser.parse_args()

    target_url   = args.target.rstrip('/')
    existing_url = args.existing.rstrip('/')
    base_path    = get_staging_base(target_url)

    if base_path:
        print(f'ステージングURL検出: ベースパス = {base_path}')

    # LLMクライアント初期化
    llm_client = None
    if args.agent:
        try:
            import anthropic
            llm_client = anthropic.Anthropic()
            print(f'LLMエージェントモード: {args.model}')
        except ImportError:
            print('anthropic パッケージが見つかりません。pip install anthropic', file=sys.stderr)
            sys.exit(1)

    # クロール
    print(f'\n構築サイトの巡回を開始: {target_url}')
    crawler = Crawler(target_url, base_path)
    crawler.crawl()
    print(f'巡回完了: {len(crawler.pages)} ページ\n')

    # 既存サイト取得
    existing_soup = None
    print(f'既存サイトのトップページを取得: {existing_url}')
    try:
        sess = requests.Session()
        sess.headers['User-Agent'] = 'Mozilla/5.0 (compatible; WebInspector/1.0)'
        resp = sess.get(existing_url, timeout=15)
        existing_soup = BeautifulSoup(resp.text, 'lxml')
        print('取得完了\n')
    except Exception as e:
        print(f'既存サイト取得エラー: {e}\n', file=sys.stderr)

    # チェック実行
    all_results: dict = {}
    for url, page_data in crawler.pages.items():
        results = run_checks(url, page_data, base_path, existing_soup, llm_client, args.model)
        all_results[url] = results
        path  = rel_path(url, base_path)
        ng_n  = sum(1 for r in results if r['kind'] == 'NG')
        w_n   = sum(1 for r in results if r['kind'] == '要手動確認')
        print(f'チェック完了: {path} (NG:{ng_n}件, 要確認:{w_n}件)')

    # ワイヤーフレーム照合
    if args.wire:
        print(f'\nワイヤーフレーム照合: {args.wire}')
        wire_results = check_wireframe(args.wire, crawler.pages)
        top_url = target_url + '/'
        key = top_url if top_url in all_results else target_url
        all_results.setdefault(key, [])
        all_results[key] = wire_results + all_results[key]

    # Excel出力
    ts  = datetime.now().strftime('%Y%m%d_%H%M%S')
    out = f'検品結果_{ts}.xlsx'
    print(f'\n結果を出力中: {out}')
    write_excel(all_results, out)

    total_ng = sum(sum(1 for r in v if r['kind'] == 'NG')        for v in all_results.values())
    total_w  = sum(sum(1 for r in v if r['kind'] == '要手動確認') for v in all_results.values())
    print(f'\n完了  NG={total_ng}件  要手動確認={total_w}件')
    print(f'出力ファイル: {out}')


if __name__ == '__main__':
    main()
